import os
import sys
import logging
from pptx import Presentation
from pptx.util import Inches
import fitz  # PyMuPDF
import tempfile
import aspose.slides as slides

logging.basicConfig(level=logging.DEBUG, format='%(asctime)s - %(name)s - %(levelname)s - %(message)s')
logger = logging.getLogger(__name__)

def convert_pptx_to_pdf(input_path: str, output_path: str):
    """
    Converts a PPTX file to PDF using aspose.slides (cross-platform, cloud-friendly).
    """
    input_path = os.path.abspath(input_path)
    output_path = os.path.abspath(output_path)
    
    try:
        logger.info(f"Starting PPTX to PDF conversion using aspose.slides. Input: {input_path}")
        
        # Load the presentation
        with slides.Presentation(input_path) as presentation:
            # Save the presentation to PDF
            presentation.save(output_path, slides.export.SaveFormat.PDF)
        
        if not os.path.exists(output_path):
            raise FileNotFoundError("PPTX to PDF conversion failed.")
        logger.info("Conversion completed successfully.")
    except Exception as e:
        import traceback
        traceback.print_exc()
        logger.exception(f"Exception during PPTX to PDF conversion: {str(e)}")
        raise RuntimeError(f"aspose.slides error: {str(e)}")


def convert_pdf_to_pptx(input_path: str, output_path: str):
    """
    Converts a PDF file to PPTX using PyMuPDF to extract images and python-pptx to assemble slides.
    """
    input_path = os.path.abspath(input_path)
    output_path = os.path.abspath(output_path)
    
    try:
        logger.info(f"Starting PDF to PPTX conversion. Input: {input_path}, Output: {output_path}")
        prs = Presentation()
        # Remove default empty slide
        xml_slides = prs.slides._sldIdLst  
        slides = list(xml_slides)
        for slide in slides:
            xml_slides.remove(slide)
            
        with fitz.open(input_path) as doc:
            if doc.needs_pass:
                doc.authenticate('')
                if doc.needs_pass:
                    raise RuntimeError("PDF is encrypted and cannot be processed.")
            
            with tempfile.TemporaryDirectory() as temp_dir:
                for page_num in range(len(doc)):
                    page = doc.load_page(page_num)
                    zoom = 2.0  # High quality
                    mat = fitz.Matrix(zoom, zoom)
                    pix = page.get_pixmap(matrix=mat, alpha=False)
                    
                    img_path = os.path.join(temp_dir, f"page_{page_num}.png")
                    pix.save(img_path)
                    
                    # Create a blank slide layout
                    blank_slide_layout = prs.slide_layouts[6]
                    slide = prs.slides.add_slide(blank_slide_layout)
                    
                    # Add picture to fill the slide
                    left = top = Inches(0)
                    width = prs.slide_width
                    height = prs.slide_height
                    slide.shapes.add_picture(img_path, left, top, width, height)
                    
        prs.save(output_path)
        
        if not os.path.exists(output_path):
            logger.error("Output file not found after conversion.")
            raise FileNotFoundError("PDF to PPTX conversion failed.")
        logger.info("Conversion completed successfully.")
    except Exception as e:
        import traceback
        traceback.print_exc()
        logger.exception(f"Exception during PDF to PPTX conversion: {str(e)}")
        raise RuntimeError(f"PyMuPDF/python-pptx error: {str(e)}")
